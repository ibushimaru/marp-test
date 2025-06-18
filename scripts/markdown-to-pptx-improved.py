#!/usr/bin/env python3
"""
改良版：Markdownから編集可能なPowerPointを生成するスクリプト
レイアウトとデザインを改善
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def parse_markdown(file_path):
    """Markdownファイルを解析してスライドデータを抽出"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # フロントマターをスキップ
    if content.startswith('---'):
        _, _, content = content.split('---', 2)
    
    # スライドを分割
    slides = content.split('\n---\n')
    
    parsed_slides = []
    for slide in slides:
        lines = slide.strip().split('\n')
        if not lines:
            continue
            
        slide_data = {
            'title': '',
            'content': [],
            'subtitle': '',
            'is_title_slide': False,
            'tables': []
        }
        
        in_table = False
        table_data = []
        
        for line in lines:
            line = line.strip()
            
            # HTMLコメントでクラス指定を確認
            if '<!-- _class: title -->' in line:
                slide_data['is_title_slide'] = True
                continue
            
            # HTMLコメントやクラス指定をスキップ
            if line.startswith('<!--') or line == '' and not in_table:
                continue
                
            # テーブルの処理
            if line.startswith('|'):
                if not in_table:
                    in_table = True
                    table_data = []
                table_data.append(line)
                continue
            elif in_table:
                if line and not line.startswith('|'):
                    # テーブル終了
                    slide_data['tables'].append(table_data)
                    in_table = False
                    table_data = []
            
            # タイトル（# または ##）
            if line.startswith('# '):
                title_text = line[2:].strip()
                # {.important}などのクラス指定を除去
                title_text = re.sub(r'{\..*?}', '', title_text)
                slide_data['title'] = title_text
            elif line.startswith('## '):
                title_text = line[3:].strip()
                title_text = re.sub(r'{\..*?}', '', title_text)
                if not slide_data['title']:
                    slide_data['title'] = title_text
                else:
                    slide_data['subtitle'] = title_text
            elif line.startswith('### '):
                text = line[4:].strip()
                text = re.sub(r'{\..*?}', '', text)
                slide_data['content'].append(('heading3', text))
            elif line.startswith('**') and line.endswith('**') and not '**' in line[2:-2]:
                # 太字のみの行
                text = line[2:-2].strip()
                text = re.sub(r'{\..*?}', '', text)
                slide_data['content'].append(('bold', text))
            elif line.startswith('- ') or line.startswith('* '):
                text = line[2:].strip()
                text = re.sub(r'{\..*?}', '', text)
                slide_data['content'].append(('bullet', text))
            elif line.startswith('1. ') or re.match(r'^\d+\.\s', line):
                text = re.sub(r'^\d+\.\s*', '', line).strip()
                text = re.sub(r'{\..*?}', '', text)
                slide_data['content'].append(('numbered', text))
            elif line and not line.startswith('```'):
                # 通常のテキスト
                text = re.sub(r'{\..*?}', '', line)
                if text:
                    slide_data['content'].append(('text', text))
        
        # 最後のテーブルを処理
        if in_table and table_data:
            slide_data['tables'].append(table_data)
        
        if slide_data['title'] or slide_data['content']:
            parsed_slides.append(slide_data)
    
    return parsed_slides

def create_pptx_improved(slides_data, output_path):
    """改良版：スライドデータからPowerPointファイルを生成"""
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # カラー定義（minimalTheme.cssから）
    color_primary = RGBColor(0x2C, 0x3E, 0x50)
    color_secondary = RGBColor(0x7F, 0x8C, 0x8D)
    color_accent_orange = RGBColor(0xFF, 0x6B, 0x35)
    color_accent_blue = RGBColor(0x34, 0x98, 0xDB)
    
    for slide_data in slides_data:
        # スライドレイアウトを選択
        if slide_data['is_title_slide']:
            slide_layout = prs.slide_layouts[0]  # Title Slide
        elif slide_data['tables']:
            slide_layout = prs.slide_layouts[5]  # Blank
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトルを設定
        if slide_data['title']:
            if slide_data['is_title_slide']:
                # タイトルスライドの場合
                title = slide.shapes.title
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.name = 'FORM UDPGothic'
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = color_primary
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # サブタイトルがある場合
                if slide_data['subtitle'] and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
                    subtitle = slide.placeholders[1]
                    subtitle.text = slide_data['subtitle']
                    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
                    subtitle.text_frame.paragraphs[0].font.name = 'FORM UDPGothic'
                    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                # 通常スライドの場合
                if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                    title = slide.shapes.title
                    title.text = slide_data['title']
                    title.text_frame.paragraphs[0].font.size = Pt(32)
                    title.text_frame.paragraphs[0].font.name = 'FORM UDPGothic'
                    title.text_frame.paragraphs[0].font.bold = True
                    title.text_frame.paragraphs[0].font.color.rgb = color_primary
                else:
                    # タイトルプレースホルダーがない場合は手動で追加
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = slide_data['title']
                    tf.paragraphs[0].font.size = Pt(32)
                    tf.paragraphs[0].font.name = 'FORM UDPGothic'
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.color.rgb = color_primary
        
        # コンテンツを設定
        if slide_data['content'] and not slide_data['tables']:
            # 通常のコンテンツスライド
            if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()
            else:
                # コンテンツプレースホルダーがない場合は手動で追加
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(3.5)
                content_shape = slide.shapes.add_textbox(left, top, width, height)
                tf = content_shape.text_frame
            
            # 行間を設定
            tf.word_wrap = True
            
            for i, (content_type, text) in enumerate(slide_data['content']):
                # マークダウン記法を除去
                clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
                clean_text = re.sub(r'{\..*?}', '', clean_text)
                
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = clean_text
                else:
                    p = tf.add_paragraph()
                    p.text = clean_text
                
                # コンテンツタイプに応じてスタイルを適用
                if content_type == 'heading3':
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.font.name = 'FORM UDPGothic'
                    p.font.color.rgb = color_primary
                    p.space_before = Pt(12)
                    p.space_after = Pt(6)
                elif content_type == 'bold':
                    p.font.bold = True
                    p.font.size = Pt(18)
                    p.font.name = 'FORM UDPGothic'
                    p.font.color.rgb = color_accent_orange
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                elif content_type == 'bullet':
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.name = 'FORM UDPGothic'
                    p.space_before = Pt(3)
                    p.space_after = Pt(3)
                elif content_type == 'numbered':
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.name = 'FORM UDPGothic'
                    p.space_before = Pt(3)
                    p.space_after = Pt(3)
                else:
                    p.font.size = Pt(16)
                    p.font.name = 'FORM UDPGothic'
                    p.space_before = Pt(3)
                    p.space_after = Pt(3)
        
        # テーブルがある場合
        if slide_data['tables']:
            for table_idx, table_data in enumerate(slide_data['tables']):
                # テーブルデータを解析
                rows = []
                for row in table_data:
                    if '|' in row and not row.strip().startswith('|-'):
                        cells = [cell.strip() for cell in row.split('|')[1:-1]]
                        rows.append(cells)
                
                if rows:
                    # テーブルを追加
                    rows_count = len(rows)
                    cols_count = len(rows[0]) if rows else 0
                    
                    left = Inches(0.5)
                    top = Inches(2 + table_idx * 2)
                    width = Inches(9)
                    height = Inches(0.5 * rows_count)
                    
                    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
                    
                    # テーブルスタイルを設定
                    for row_idx, row_data in enumerate(rows):
                        for col_idx, cell_text in enumerate(row_data):
                            cell = table.cell(row_idx, col_idx)
                            cell.text = cell_text
                            
                            # フォント設定
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.name = 'FORM UDPGothic'
                                paragraph.font.size = Pt(14)
                                
                                # ヘッダー行（最初の行）は太字
                                if row_idx == 0:
                                    paragraph.font.bold = True
                                    paragraph.font.color.rgb = color_primary
                                else:
                                    paragraph.font.color.rgb = color_secondary
    
    # ファイルを保存
    prs.save(output_path)

def main():
    input_file = '/root/claude-docker-dev/workspace/Obsidian/awesome-marp-template/slides/generative-ai-translation-minimal/index.md'
    output_file = '/root/claude-docker-dev/workspace/Obsidian/awesome-marp-template/slides/generative-ai-translation-minimal/presentation-improved.pptx'
    
    print(f"Markdownファイルを読み込み中: {input_file}")
    slides_data = parse_markdown(input_file)
    
    print(f"{len(slides_data)}枚のスライドを検出")
    
    print(f"PowerPointファイルを生成中: {output_file}")
    create_pptx_improved(slides_data, output_file)
    
    print("完了！")

if __name__ == '__main__':
    main()