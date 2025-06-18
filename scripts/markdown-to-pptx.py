#!/usr/bin/env python3
"""
Markdownから編集可能なPowerPointを生成するスクリプト
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def parse_markdown(file_path):
    """Markdownファイルを解析してスライドデータを抽出"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # フロントマターをスキップ
    if content.startswith('---'):
        _, _, content = content.split('---', 2)
    
    # スライドを分割
    slides = content.split('\n---\n')
    
    parsed_slides = []
    for slide in slides:
        lines = slide.strip().split('\n')
        if not lines:
            continue
            
        slide_data = {
            'title': '',
            'content': [],
            'subtitle': ''
        }
        
        for line in lines:
            line = line.strip()
            
            # HTMLコメントやクラス指定をスキップ
            if line.startswith('<!--') or '{.' in line:
                continue
                
            # タイトル（# または ##）
            if line.startswith('# '):
                slide_data['title'] = line[2:].strip()
            elif line.startswith('## '):
                if not slide_data['title']:
                    slide_data['title'] = line[3:].strip()
                else:
                    slide_data['subtitle'] = line[3:].strip()
            elif line.startswith('### '):
                slide_data['content'].append(('heading3', line[4:].strip()))
            elif line.startswith('**') and line.endswith('**'):
                # 太字のみの行
                slide_data['content'].append(('bold', line[2:-2].strip()))
            elif line.startswith('- ') or line.startswith('* '):
                slide_data['content'].append(('bullet', line[2:].strip()))
            elif line.startswith('1. '):
                slide_data['content'].append(('numbered', line[3:].strip()))
            elif line.startswith('|'):
                # テーブル（簡易処理）
                slide_data['content'].append(('table', line))
            elif line:
                slide_data['content'].append(('text', line))
        
        if slide_data['title'] or slide_data['content']:
            parsed_slides.append(slide_data)
    
    return parsed_slides

def create_pptx(slides_data, output_path):
    """スライドデータからPowerPointファイルを生成"""
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for slide_data in slides_data:
        # タイトルスライドかコンテンツスライドかを判断
        if not slide_data['content'] and slide_data['title']:
            slide_layout = prs.slide_layouts[0]  # Title Slide
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトルを設定
        if slide_data['title']:
            title = slide.shapes.title
            title.text = slide_data['title']
            # タイトルのフォントサイズとフォントファミリーを調整
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.name = 'FORM UDPGothic'
        
        # サブタイトルまたはコンテンツを設定
        if slide_data['subtitle'] and not slide_data['content']:
            # タイトルスライドのサブタイトル
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data['subtitle']
        elif slide_data['content']:
            # コンテンツスライド
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()  # 既存のテキストをクリア
            
            for i, (content_type, text) in enumerate(slide_data['content']):
                # クリーンなテキストを取得（マークダウン記法を除去）
                clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
                clean_text = re.sub(r'{\..*?}', '', clean_text)
                
                if i == 0:
                    tf.text = clean_text
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    p.text = clean_text
                
                # コンテンツタイプに応じてスタイルを適用
                if content_type == 'heading3':
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.name = 'FORM UDPGothic'
                elif content_type == 'bold':
                    p.font.bold = True
                    p.font.size = Pt(20)
                    p.font.name = 'FORM UDPGothic'
                elif content_type in ['bullet', 'numbered']:
                    p.level = 1
                    p.font.size = Pt(18)
                    p.font.name = 'FORM UDPGothic'
                else:
                    p.font.size = Pt(18)
                    p.font.name = 'FORM UDPGothic'
    
    # ファイルを保存
    prs.save(output_path)

def main():
    input_file = '/root/claude-docker-dev/workspace/Obsidian/awesome-marp-template/slides/generative-ai-translation-minimal/index.md'
    output_file = '/root/claude-docker-dev/workspace/Obsidian/awesome-marp-template/slides/generative-ai-translation-minimal/presentation-python.pptx'
    
    print(f"Markdownファイルを読み込み中: {input_file}")
    slides_data = parse_markdown(input_file)
    
    print(f"{len(slides_data)}枚のスライドを検出")
    
    print(f"PowerPointファイルを生成中: {output_file}")
    create_pptx(slides_data, output_file)
    
    print("完了！")

if __name__ == '__main__':
    main()